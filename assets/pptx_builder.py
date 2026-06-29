"""PptxBuilder — same interface as DrawioBuilder, but renders NATIVE PowerPoint
shapes (rounded rectangles, text boxes, ovals, connectors) + AWS icons as
individual images. Every element is an editable PPT object, not a flat picture.

    from pptx_builder import PptxBuilder
    d = PptxBuilder(w=1672, h=945, icon_dir="/tmp/icons", shadow_ids={"g_a"})
    layout(d)              # the SAME layout(d) you wrote for DrawioBuilder
    d.save("out.pptx")

Icons: run crop_icons.py first to populate icon_dir with <alias>.png. Any AWS
alias without a cropped PNG falls back to a coloured tile.
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from aws_icons import AWS4, CAT, soft

_ALN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_ANC = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


class PptxBuilder:
    def __init__(self, w=1672, h=945, icon_dir=None, shadow_ids=(), slide_in=13.333):
        self.W, self.H = w, h
        self.icon_dir = icon_dir
        self.shadow = set(shadow_ids)
        self.GEO = {}
        self.P = Presentation()
        self.EMUPX = int(slide_in * 914400) / w
        self.P.slide_width = Emu(int(slide_in * 914400))
        self.P.slide_height = Emu(int(round(h * self.EMUPX)))
        self.slide = self.P.slides.add_slide(self.P.slide_layouts[6])
        self.shapes = self.slide.shapes
        self.SF = self.EMUPX / 12700.0  # px -> pt

    # ---- low level ----
    def _E(self, px): return Emu(int(round(px * self.EMUPX)))
    def _FP(self, px): return Pt(max(4.0, round(px * self.SF, 1)))
    def _C(self, h): return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))
    def _noshadow(self, sp):
        try: sp.shadow.inherit = False
        except Exception: pass
    def _shadow(self, sp):
        spPr = sp._element.spPr
        e = spPr.find(qn('a:effectLst'))
        if e is not None: spPr.remove(e)
        el = spPr.makeelement(qn('a:effectLst'), {})
        sh = el.makeelement(qn('a:outerShdw'), {'blurRad': '50000', 'dist': '25000', 'dir': '2700000', 'rotWithShape': '0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '8A97A2'})
        a = clr.makeelement(qn('a:alpha'), {'val': '30000'}); clr.append(a); sh.append(clr); el.append(sh); spPr.append(el)
    def _text(self, tf, label, font, fc, bold, align, valign):
        tf.word_wrap = True; tf.vertical_anchor = _ANC[valign]
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, Emu(6000))
        first = True
        for ln in str(label).split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.alignment = _ALN[align]
            r = p.add_run(); r.text = ln
            r.font.size = self._FP(font); r.font.bold = bool(bold); r.font.color.rgb = self._C(fc)
    def _rect(self, id, label, x, y, w, h, fill, stroke, font, fc, bold, align, valign, arc, shadow):
        sp = self.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self._E(x), self._E(y), self._E(w), self._E(h))
        try: sp.adjustments[0] = arc
        except Exception: pass
        if fill == "none": sp.fill.background()
        else: sp.fill.solid(); sp.fill.fore_color.rgb = self._C(fill)
        if stroke == "none": sp.line.fill.background()
        else: sp.line.color.rgb = self._C(stroke); sp.line.width = Pt(0.5 if arc <= 0.1 else 0.6)
        self._noshadow(sp)
        if shadow: self._shadow(sp)
        self._text(sp.text_frame, label, font, fc, bold, align, valign)
        return sp
    def _tb(self, x, y, w, h, label, font, fc, bold, align, valign, fill=None):
        tb = self.shapes.add_textbox(self._E(x), self._E(y), self._E(w), self._E(h))
        if fill: tb.fill.solid(); tb.fill.fore_color.rgb = self._C(fill); tb.line.fill.background()
        self._text(tb.text_frame, label, font, fc, bold, align, valign)
        return tb

    # ---- public API (mirrors DrawioBuilder) ----
    def box(self, id, label, x, y, w, h, fill="#FFFFFF", stroke="#2D3748", font=8,
            fc="#1A202C", bold=0, align="center", valign="middle"):
        self._rect(id, label, x, y, w, h, fill, soft(stroke, 0.45) if stroke != "none" else "none",
                   font, fc, bold, align, valign, 0.08, False)
        if id: self.GEO[id] = (x, y, w, h)

    def grp(self, id, label, x, y, w, h, stroke, fill="none", font=9):
        self._rect(id, label, x, y, w, h, fill, soft(stroke, 0.3), font, stroke, 1, "center", "top", 0.04, id in self.shadow)
        self.GEO[id] = (x, y, w, h)

    def band(self, id, label, x, y, w, h, c, font=10, fc="#FFFFFF"):
        self._rect(id, label, x, y, w, h, c, "none", font, fc, 1, "center", "middle", 0.16, False)
        if id: self.GEO[id] = (x, y, w, h)

    def txt(self, id, label, x, y, w, h=16, font=10, bold=1, align="left", fc="#1A202C"):
        self._tb(x, y, w, h, label, font, fc, bold, align, "middle")
        if id: self.GEO[id] = (x, y, w, h)

    def aws(self, id, label, x, y, svc, sz=32, catov=None):
        if svc not in AWS4:
            return self.box(id, label, x, y, sz, sz, "#FFFFFF", "#FF9900", 7)
        res, cat = AWS4[svc]; cat = catov or cat; fill, _ = CAT[cat]
        p = os.path.join(self.icon_dir, svc + ".png") if self.icon_dir else None
        if p and os.path.exists(p):
            self.shapes.add_picture(p, self._E(x), self._E(y), self._E(sz), self._E(sz))
        else:
            self.box(None, "", x, y, sz, sz, fill, "#FFFFFF", 6)
        if label:
            lw = sz * 2.8
            self._tb(x + sz / 2 - lw / 2, y + sz + 1, lw, 16, label, 6, "#232F3E", 0, "center", "top")
        if id: self.GEO[id] = (x, y, sz, sz)

    def emoji(self, glyph, cx, cw, y, sz):
        self._tb(round(cx + cw / 2 - sz / 2), y, sz, sz, glyph, sz - 8, "#232F3E", 0, "center", "middle")

    def ic_tb(self, idp, svc, lab, cx, cw, y, sz=32, catov=None):
        if svc in AWS4: self.aws(idp, "", round(cx + cw / 2 - sz / 2), y, svc, sz, catov)
        else: self.emoji(svc, cx, cw, y, sz)
        self.box(idp + "t", lab, round(cx), y + sz + 1, round(cw), 22, "none", "none", 6, align="center", valign="top")

    def ellipse(self, id, glyph, x, y, w, h, stroke, fc="#111111", font=14):
        el = self.shapes.add_shape(MSO_SHAPE.OVAL, self._E(x), self._E(y), self._E(w), self._E(h))
        el.fill.solid(); el.fill.fore_color.rgb = self._C("#FFFFFF")
        el.line.color.rgb = self._C(stroke); el.line.width = Pt(1.5); self._noshadow(el)
        self._text(el.text_frame, glyph, font, fc, 0, "center", "middle")

    def pcard(self, idp, emo, title, desc, x, y, w, h, c):
        self.box(idp, "", x, y, w, h, "#FFFFFF", c, 7)
        self.ellipse(idp + "b", emo, x + 8, round(y + h / 2 - 15), 30, 30, c)
        self.box(idp + "tt", title, x + 46, y + 6, w - 50, 16, "none", "none", 8, fc="#1A202C", bold=1, align="left", valign="middle")
        self.box(idp + "td", desc, x + 46, y + 22, w - 50, h - 26, "none", "none", 7, fc="#5B6B7B", align="left", valign="top")

    def bv(self, idp, glyph, lab, cx, cw, y, color):
        self.ellipse(idp, glyph, round(cx + cw / 2 - 15), y, 30, 30, color, fc=color)
        self.box(idp + "t", lab, round(cx), y + 32, round(cw), 24, "none", "none", 6, align="center", valign="top")

    def edge(self, s, d, label="", dashed=False):
        if s not in self.GEO or d not in self.GEO: return
        sx, sy, sw, sh = self.GEO[s]; dx, dy, dw, dh = self.GEO[d]
        if dy > sy + sh - 5: x1, y1, x2, y2 = sx + sw / 2, sy + sh, dx + dw / 2, dy
        elif dx >= sx + sw - 5: x1, y1, x2, y2 = sx + sw, sy + sh / 2, dx, dy + dh / 2
        else: x1, y1, x2, y2 = sx + sw / 2, sy + sh, dx + dw / 2, dy
        cn = self.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, self._E(x1), self._E(y1), self._E(x2), self._E(y2))
        cn.line.color.rgb = self._C("#6B7B8C"); cn.line.width = Pt(1.1); self._noshadow(cn)
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        if dashed: ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
        if label:
            mx, my = (x1 + x2) / 2, (y1 + y2) / 2
            self._tb(mx - 32, my - 8, 64, 14, label, 7, "#37474F", 1, "center", "middle", fill="#FFFFFF")

    def save(self, path):
        self.P.save(path)
        return path
